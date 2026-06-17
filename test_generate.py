"""Verificación local de la generación del deck SIN token de Notion.

Usa un fixture representativo (reel, carrusel, historia y borrador) para
confirmar que slides.build_deck produce un .pptx válido con el diseño esperado.

    python test_generate.py
"""

from pptx import Presentation

from lib.slides import build_deck

FIXTURE = [
    dict(
        date="01 Jul", name="Historia 1", icon="⏰",
        formato="Historias en imagenes", red="Instagram Stories",
        pilar="Contenido de venta", tipo="Venta directa",
        objetivo="Conversion", estado="En revision", borrador=False,
        guion=[
            "H1: Comienza la semana de la dulzura y tenemos el detalle perfecto para regalar",
            "H2: Mostrar productos mini con flechas: facturas / alfajores / cookies",
            "H3: Caja de bombones + flores para tu pareja",
            "H4: No te quedes sin hacer un regalo (link WhatsApp)",
        ],
        copy=None,
    ),
    dict(
        date="02 Jul", name="Reel 1", icon="\U0001f4fd",
        formato="Reel", red="Instagram",
        pilar="Contenido de venta indirecta", tipo=None,
        objetivo="Alcance / Interaccion", estado="En revision", borrador=False,
        guion=[
            "Hook: Mira estas opciones para regalar por la semana de la dulzura",
            "Desarrollo: opciones mini + caja de bombones con flores",
            "Cierre: pasa a buscar tu mejor regalo en La Portena",
        ],
        copy="La semana de la dulzura esta con nosotros! Cuales son tu opcion favorita? Te leemos en los comentarios.",
    ),
    dict(
        date="03 Jul", name="Historia 2 libre", icon="⏰",
        formato="Historias en imagenes", red="Instagram Stories",
        pilar=None, tipo=None, objetivo=None,
        estado="Borrador", borrador=True, guion=None, copy=None,
    ),
    dict(
        date="05 Jul", name="Posteo 4", icon="\U0001f5bc",
        formato="Carrusel", red="Instagram Feed",
        pilar="Contenido de valor", tipo=None,
        objetivo="Alcance / Interaccion", estado="En revision", borrador=False,
        guion=[
            "Img 1: Unas fotos no te pueden dar tanto hambre, las fotos:",
            "Img 2: Foto de croissant de pistacho",
            "Img 3: Foto de torta matilda",
            "Img 4: Compartile este posteo a esa amiga fan de lo dulce",
        ],
        copy="Levante la mano quien ve algo dulce y le agarra hambre! Compartile a esa amiga. #panaderia #moreno",
    ),
]


def main():
    buf = build_deck(FIXTURE, "Julio  2026")
    with open("test_output.pptx", "wb") as f:
        f.write(buf.getvalue())

    prs = Presentation("test_output.pptx")
    n = len(prs.slides)
    expected = len(FIXTURE) + 1  # +1 portada
    assert n == expected, f"Esperaba {expected} slides, hay {n}"

    # Chequeo de textos clave
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    blob = "\n".join(all_text)
    assert "Planificacion de Contenidos" in blob
    assert "Reel 1" in blob
    assert "CONTENIDO PENDIENTE" in blob
    assert "GUION" in blob and "COPY" in blob

    print(f"OK · {n} slides generadas · test_output.pptx")
    for i, slide in enumerate(prs.slides):
        first = next((s.text_frame.text for s in slide.shapes
                      if s.has_text_frame and s.text_frame.text.strip()), "")
        print(f"  Slide {i+1}: {first[:50]}")


if __name__ == "__main__":
    main()
