"""Generador de presentaciones .pptx — identidad de marca MATTER.

Estética: minimal clara. Fondo crema, texto negro, acentos lima y lavanda.
Títulos en Archivo Black (similar a Akira Expanded), cuerpo en Inter.

Función principal:
    build_deck(items, deck_title, client_name=None) -> io.BytesIO

Cada `item` es un dict con:
    date, name, icon, red, estado, formato, pilar, tipo, objetivo,
    borrador (bool), guion (list[str] | None), copy (str | None)
"""

import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Paleta MATTER ────────────────────────────────────────────────────────────
BLACK   = RGBColor(0x00, 0x00, 0x00)
CREAM   = RGBColor(0xF1, 0xE9, 0xDE)
LIME    = RGBColor(0xD2, 0xFF, 0x70)
LAVENDER = RGBColor(0xCB, 0xAC, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x8A, 0x82, 0x77)   # gris cálido sobre crema
HAIRLINE = RGBColor(0x00, 0x00, 0x00)
DRAFT_MUTE = RGBColor(0xB4, 0xAC, 0xA1)

# Tipografías (disponibles en Google Slides).
F_TITLE = "Archivo Black"
F_BODY  = "Inter"

# Dimensiones (pulgadas).
PAGE_W = 10.0
PAGE_H = 5.625
MARGIN = 0.6


# ─── Helpers de texto ─────────────────────────────────────────────────────────
def _set_spacing(run, pts):
    """Letter-spacing en puntos (para el look 'expanded' de los títulos)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(pts * 100)))


def _add_text(slide, text, x, y, w, h, size, *, bold=False, italic=False,
              color=BLACK, align=PP_ALIGN.LEFT, font=F_BODY, spacing=None,
              anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    if spacing is not None:
        _set_spacing(r, spacing)
    return tb


def _add_multiline(slide, lines, x, y, w, h, size, color, *, font=F_BODY,
                   space_after=4, line_spacing=1.04):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0)
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
    return tb


def _add_rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _add_oval(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _add_line(slide, x, y, w, color=HAIRLINE, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def _pill_w(text, size=10.5):
    """Ancho estimado de una pastilla para un texto y tamaño dados."""
    factor = size / 10.5
    return max(0.82, (0.082 * len(text) + 0.42) * factor)


def _add_pill(slide, x, y, text, fill, *, text_color=BLACK, size=10.5,
              font=F_BODY, bold=True):
    """Pastilla redondeada con texto centrado. Devuelve el ancho usado."""
    w = _pill_w(text, size)
    h = 0.30
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    r.font.name = font
    return w


# ─── Portada ──────────────────────────────────────────────────────────────────
def _cover(slide, deck_title, client_name, n_items):
    _add_rect(slide, 0, 0, PAGE_W, PAGE_H, CREAM)

    # Motivo de marca: cuadrado lima + wordmark MATTER (arriba-izq).
    _add_rect(slide, MARGIN, 0.5, 0.26, 0.26, LIME)
    _add_text(slide, "MATTER", MARGIN + 0.4, 0.46, 4.0, 0.4, 17,
              bold=True, color=BLACK, font=F_TITLE, spacing=1.5)

    # Motivos geométricos minimal (esquina sup. der).
    _add_oval(slide, 8.9, 0.5, 0.5, 0.5, LAVENDER)
    _add_rect(slide, 8.55, 0.62, 0.26, 0.26, LIME)

    # Bloque central.
    _add_text(slide, "PLANIFICACIÓN DE CONTENIDOS", MARGIN, 1.95, 8.8, 0.4, 13,
              bold=True, color=MUTED, font=F_BODY, spacing=2.5)

    big = (client_name or deck_title or "Contenido").upper()
    _add_text(slide, big, MARGIN, 2.35, 8.8, 1.4, 46,
              bold=True, color=BLACK, font=F_TITLE, spacing=0.5)

    # Barra acento lima debajo del título.
    _add_rect(slide, MARGIN + 0.02, 3.55, 2.2, 0.12, LIME)

    # Si hay cliente, el deck_title (mes) va como subtítulo.
    if client_name and deck_title:
        _add_text(slide, deck_title, MARGIN, 3.78, 8.8, 0.5, 18,
                  bold=False, color=BLACK, font=F_BODY)

    # Pie de portada.
    _add_text(slide, f"{n_items} piezas de contenido", MARGIN, 4.95, 6.0, 0.4, 12,
              color=MUTED, font=F_BODY)


# ─── Slide de contenido ─────────────────────────────────────────────────────
def _date_right(slide, date_txt):
    """Línea de fecha alineada a la derecha: 'PUBLICACIÓN  05 Jul'."""
    tb = slide.shapes.add_textbox(Inches(PAGE_W - MARGIN - 3.4), Inches(0.84),
                                  Inches(3.4), Inches(0.32))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run()
    r1.text = "PUBLICACIÓN   "
    r1.font.size = Pt(9)
    r1.font.bold = True
    r1.font.color.rgb = MUTED
    r1.font.name = F_BODY
    _set_spacing(r1, 1.2)
    r2 = p.add_run()
    r2.text = date_txt
    r2.font.size = Pt(15)
    r2.font.bold = True
    r2.font.color.rgb = BLACK
    r2.font.name = F_BODY


def _content_slide(slide, item):
    is_draft = item.get("borrador", False)
    _add_rect(slide, 0, 0, PAGE_W, PAGE_H, CREAM)
    # Spine de color a la izquierda (lima; lavanda si es borrador-distinto).
    _add_rect(slide, 0, 0, 0.18, PAGE_H, DRAFT_MUTE if is_draft else LIME)

    # ── Encabezado: nombre de la pieza (izquierda, más chico) ────────────────
    name = (item.get("name") or "Sin nombre").upper()
    _add_text(slide, name, MARGIN, 0.46, 5.0, 0.6, 24,
              bold=True, color=BLACK, font=F_TITLE, spacing=0.3)

    # ── Lado derecho: pastilla red social (ícono chico) + fecha ──────────────
    icon = item.get("icon") or ""
    red = item.get("red") or item.get("formato") or "Instagram"
    red_label = f"{icon} {red}".strip()
    pw = _pill_w(red_label, 9.5)
    _add_pill(slide, PAGE_W - MARGIN - pw, 0.44, red_label, LIME, size=9.5)

    date_txt = item.get("date") or ""
    if date_txt:
        _date_right(slide, date_txt)

    # Hairline separador.
    _add_line(slide, MARGIN, 1.3, PAGE_W - 2 * MARGIN, color=BLACK, weight=1.0)

    if is_draft:
        _add_text(slide, "CONTENIDO PENDIENTE", MARGIN, 2.7, PAGE_W - 2 * MARGIN, 0.8,
                  26, bold=True, color=DRAFT_MUTE, font=F_TITLE,
                  align=PP_ALIGN.CENTER, spacing=0.5)
        _add_text(slide, "Guión pendiente de elaboración", MARGIN, 3.5,
                  PAGE_W - 2 * MARGIN, 0.4, 12, italic=True, color=DRAFT_MUTE,
                  font=F_BODY, align=PP_ALIGN.CENTER)
        return

    # ── Cuerpo (más espacio: sin footer, arranca más arriba) ─────────────────
    guion = item.get("guion") or []
    copy = item.get("copy")
    has_guion = bool(guion)
    has_copy = bool(copy)
    body_y = 1.55
    body_h = 3.85

    def _section_label(text, x, y, w_, color=BLACK):
        _add_text(slide, text, x, y, w_, 0.26, 10, bold=True, color=color,
                  font=F_BODY, spacing=1.8)

    if has_guion and has_copy:
        # Dos columnas: izquierda guión, derecha copy.
        lx, lw = MARGIN, 4.55
        rx, rw = 5.5, 3.9
        # Divisor vertical sutil.
        vdiv = slide.shapes.add_connector(2, Inches(5.18), Inches(body_y),
                                          Inches(5.18), Inches(body_y + body_h))
        vdiv.line.color.rgb = RGBColor(0xD8, 0xCE, 0xC0)
        vdiv.line.width = Pt(0.75)
        vdiv.shadow.inherit = False

        _section_label("QUÉ VA EN CADA SLIDE", lx, body_y, lw)
        n = len(guion)
        gsize = 12 if n <= 3 else (10.5 if n <= 5 else 9.5)
        _add_multiline(slide, guion, lx, body_y + 0.34, lw, body_h - 0.34,
                       gsize, BLACK, space_after=5)

        _section_label("COPY", rx, body_y, rw)
        _add_multiline(slide, [copy], rx, body_y + 0.34, rw, body_h - 0.34,
                       11, BLACK, space_after=4, line_spacing=1.14)

    elif has_guion:
        _section_label("QUÉ VA EN CADA SLIDE", MARGIN, body_y, PAGE_W - 2 * MARGIN)
        n = len(guion)
        gsize = 14 if n <= 3 else (12 if n <= 5 else 10.5)
        _add_multiline(slide, guion, MARGIN, body_y + 0.36, PAGE_W - 2 * MARGIN,
                       body_h - 0.36, gsize, BLACK, space_after=7)

    elif has_copy:
        _section_label("COPY", MARGIN, body_y, PAGE_W - 2 * MARGIN)
        _add_multiline(slide, [copy], MARGIN, body_y + 0.36, PAGE_W - 2 * MARGIN,
                       body_h - 0.36, 14, BLACK, space_after=4, line_spacing=1.25)


# ─── Construcción del deck ─────────────────────────────────────────────────────
def build_deck(items, deck_title, client_name=None):
    prs = Presentation()
    prs.slide_width = Inches(PAGE_W)
    prs.slide_height = Inches(PAGE_H)
    blank = prs.slide_layouts[6]

    cover = prs.slides.add_slide(blank)
    _cover(cover, deck_title, client_name, len(items))

    for item in items:
        slide = prs.slides.add_slide(blank)
        _content_slide(slide, item)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
